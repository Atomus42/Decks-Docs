"""
Generate ArcaScience I-demo Q&A Committee Response Deck
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Brand Colors ───
DARK_NAVY = RGBColor(0x03, 0x17, 0x34)       # #031734 - primary dark
NAVY = RGBColor(0x0E, 0x47, 0x7A)            # #0E477A - secondary dark
BLUE = RGBColor(0x19, 0x73, 0xC2)            # #1973C2 - primary blue
CYAN = RGBColor(0x3A, 0xCF, 0xFF)            # #3ACFFF - accent cyan
AMBER = RGBColor(0xF1, 0xB4, 0x57)           # #F1B457 - accent gold
ORANGE = RGBColor(0xFF, 0xAB, 0x40)          # #FFAB40 - accent orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY = RGBColor(0x99, 0x99, 0x99)
DARK_GREY = RGBColor(0x43, 0x43, 0x43)
SOFT_BLUE_BG = RGBColor(0xE5, 0xEB, 0xF8)    # #E5EBF8

# WP color tags
WP_COLORS = {
    "WP1": RGBColor(0x19, 0x73, 0xC2),       # blue
    "WP2": RGBColor(0x00, 0x96, 0x88),        # teal
    "WP3": RGBColor(0x4C, 0xAF, 0x50),        # green
    "WP4": RGBColor(0xFF, 0xAB, 0x40),        # orange
    "WP5": RGBColor(0x7B, 0x1F, 0xA2),        # purple
    "WP6": RGBColor(0xE5, 0x39, 0x35),        # red
    "WP7": RGBColor(0x75, 0x75, 0x75),        # grey
    "Cross": RGBColor(0x2D, 0x2D, 0x2D),      # black
    "Gouvernance": RGBColor(0x2D, 0x2D, 0x2D),
    "Stratégie": RGBColor(0x2D, 0x2D, 0x2D),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "logo_header.png")

# ─── Presentation setup ───
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_dark_bg(slide):
    """Add dark navy background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY


def add_light_bg(slide):
    """Add light background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_footer(slide, page_num):
    """Add footer bar with text and page number."""
    # Footer bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_HEIGHT - Inches(0.45),
        SLIDE_WIDTH, Inches(0.45)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"ArcaScience — I-demo Q&A — Confidentiel — Avril 2026          {page_num}"
    p.font.size = Pt(9)
    p.font.color.rgb = MID_GREY
    p.alignment = PP_ALIGN.CENTER


def add_logo(slide):
    """Add logo to top right."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            SLIDE_WIDTH - Inches(1.5) - Inches(0.2),
            Inches(0.15),
            Inches(1.5),
            Inches(0.42)
        )


def add_wp_tag(slide, wp_text):
    """Add WP color tag in top right."""
    color = WP_COLORS.get(wp_text, WP_COLORS["Cross"])
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SLIDE_WIDTH - Inches(3.0),
        Inches(0.15),
        Inches(1.2),
        Inches(0.35)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    tf = tag.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = wp_text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_question_box(slide, q_fr, q_en, top=Inches(0.85)):
    """Add the question box with French + English."""
    # Box background
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top,
        SLIDE_WIDTH - Inches(1.0), Inches(1.15)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT_BLUE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = q_fr
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = DARK_NAVY

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = q_en
    run2.font.size = Pt(9)
    run2.font.italic = True
    run2.font.color.rgb = MID_GREY


def add_title_text(slide, title_text):
    """Add slide title."""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(8.5), Inches(0.65)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY


def add_content_bullets(slide, sections, left=Inches(0.6), top=Inches(2.15),
                        width=None, height=Inches(4.5), two_col=False):
    """
    Add structured bullet content.
    sections: list of (header, [bullet_texts]) or just bullet strings.
    If two_col=True, split into two columns.
    """
    if width is None:
        width = SLIDE_WIDTH - Inches(1.2) if not two_col else Inches(5.8)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)

    first = True
    for item in sections:
        if isinstance(item, tuple):
            header, bullets = item
            # Section header
            if not first:
                sp = tf.add_paragraph()
                sp.space_before = Pt(4)
                sp.space_after = Pt(0)
                sp.text = ""
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            run = p.add_run()
            run.text = header
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = BLUE
            p.space_before = Pt(4)
            p.space_after = Pt(2)
            # Bullets
            for b in bullets:
                bp = tf.add_paragraph()
                bp.level = 1
                run = bp.add_run()
                # Check for DRAFT marker
                if b.startswith("[DRAFT]"):
                    run.text = b
                    run.font.size = Pt(11)
                    run.font.color.rgb = ORANGE
                    run.font.italic = True
                else:
                    run.text = b
                    run.font.size = Pt(11)
                    run.font.color.rgb = DARK_GREY
                bp.space_before = Pt(1)
                bp.space_after = Pt(1)
        else:
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            run = p.add_run()
            if item.startswith("[DRAFT]"):
                run.text = item
                run.font.size = Pt(11)
                run.font.color.rgb = ORANGE
                run.font.italic = True
            else:
                run.text = item
                run.font.size = Pt(11)
                run.font.color.rgb = DARK_GREY
            p.space_before = Pt(2)

    return txBox


def add_draft_box(slide, text, owner, left=Inches(0.6), top=Inches(4.8),
                  width=None, height=Inches(1.8)):
    """Add a DRAFT answer box with dashed-border look."""
    if width is None:
        width = SLIDE_WIDTH - Inches(1.2)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    box.line.color.rgb = ORANGE
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"DRAFT — à valider par {owner}"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = ORANGE

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = text
    run2.font.size = Pt(10)
    run2.font.color.rgb = DARK_GREY

    return box


def add_placeholder_box(slide, label, left=Inches(0.6), top=Inches(2.15),
                         width=None, height=Inches(0.5)):
    """Add a grey placeholder box."""
    if width is None:
        width = SLIDE_WIDTH - Inches(1.2)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"À compléter — {label}"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = MID_GREY
    p.alignment = PP_ALIGN.CENTER


def make_question_slide(q_num, short_title, q_fr, q_en, sections, wp_tag,
                        speaker_notes="", page_num=0, draft_box=None,
                        placeholder=None, two_col_sections=None):
    """Create a standard question slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_light_bg(slide)
    add_title_text(slide, f"Q{q_num} — {short_title}")
    add_question_box(slide, q_fr, q_en)
    add_wp_tag(slide, wp_tag)
    add_logo(slide)
    add_footer(slide, page_num)

    if placeholder:
        add_placeholder_box(slide, placeholder, top=Inches(2.10))
        content_top = Inches(2.70)
    else:
        content_top = Inches(2.15)

    if two_col_sections:
        # Left column
        add_content_bullets(slide, sections,
                            left=Inches(0.5), top=content_top,
                            width=Inches(5.8), height=Inches(4.0))
        # Right column
        add_content_bullets(slide, two_col_sections,
                            left=Inches(6.6), top=content_top,
                            width=Inches(5.8), height=Inches(4.0))
    else:
        if draft_box:
            add_content_bullets(slide, sections,
                                top=content_top, height=Inches(2.5))
        else:
            add_content_bullets(slide, sections, top=content_top)

    if draft_box:
        d_top = Inches(4.6) if not placeholder else Inches(4.8)
        add_draft_box(slide, draft_box[0], draft_box[1], top=d_top, height=Inches(2.2))

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes

    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Title
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "I-demo — Réponses aux questions\ndu comité d'évaluation"
run.font.size = Pt(36)
run.font.bold = True
run.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "BR-PREDICT / ArcaScience"
run2.font.size = Pt(28)
run2.font.color.rgb = CYAN
run2.font.bold = False
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(3.5), Inches(4.8), Inches(6.3), Inches(1.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
run3 = p3.add_run()
run3.text = "Avril 2026 — Confidentiel"
run3.font.size = Pt(16)
run3.font.color.rgb = AMBER
p3.alignment = PP_ALIGN.CENTER

# Accent line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(4.5), Inches(4.5),
    Inches(4.3), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = CYAN
line.line.fill.background()

add_logo(slide)
add_footer(slide, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_title_text(slide, "Agenda — Thématiques abordées")
add_logo(slide)
add_footer(slide, 2)

agenda_items = [
    ("1.", "Infrastructure & Ressources", "Q1"),
    ("2.", "RH & Dissémination scientifique", "Q2–Q3"),
    ("3.", "LLM & Explicabilité", "Q4–Q5"),
    ("4.", "Trustworthy AI & Biais", "Q6–Q7"),
    ("5.", "Validation & Golden Standard", "Q8–Q10"),
    ("6.", "Monitoring & Déploiement on-premise", "Q11–Q12"),
    ("7.", "Sous-traitance & PI", "Q13–Q16"),
    ("8.", "Structure capitalistique & Financement", "Q20–Q23"),
    ("9.", "Modèle économique & Positionnement", "Q24–Q30"),
]

y = Inches(1.1)
for num, title, qs in agenda_items:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), y,
        Inches(11.3), Inches(0.55)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT_BLUE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{num}  {title}"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY
    # Question range on right
    run2 = p.add_run()
    run2.text = f"          {qs}"
    run2.font.size = Pt(12)
    run2.font.color.rgb = BLUE
    y += Inches(0.65)

# ═══════════════════════════════════════════════════════════════
# QUESTION SLIDES
# ═══════════════════════════════════════════════════════════════
page = 3

# ─── Q1: Infrastructure & Calcul ───
make_question_slide(
    q_num=1,
    short_title="Capacités de stockage et calcul",
    q_fr="De quelles capacités de stockage et de calcul dispose actuellement ArcaScience ? Quels sont les sous-traitants ? Des dépenses de cette nature sont précisées dans les WP5 (infra Neo4j cloud) et WP6 (cloud GPU haute performance), quelles ressources sont allouées aux WP1-4 ?",
    q_en="What storage and computing capabilities does ArcaScience currently have? Who are the subcontractors? What resources are allocated to WP1-4?",
    wp_tag="Cross",
    sections=[
        ("Plateforme actuelle", [
            "DataFactory BRA indexe >40M entrées textuelles (PubMed, MEDLINE, registres d'essais cliniques)",
            "Feuille de route technique définie sous T2.1, planifiée Q2–Q3 2026",
        ]),
        ("Calcul", [
            "WP1, WP2 : clusters CPU standards — pas de GPU nécessaire",
            "WP3, WP4, WP6 : environnements GPU supplémentaires, dimensionnés après validation focalisée",
            "MLflow : automatisation ML, tracking, versioning des modèles, stockage des artefacts",
        ]),
        ("Stockage", [
            "Extension infra existante : Buckets S3, ElasticSearch, Qdrant pour datasets WP-spécifiques",
            "KG WP5 : évalué sur Neo4j + backends analytiques alternatifs (benchmarks cold/cached = KPIs WP5)",
        ]),
        ("Sous-traitants", [
            "Infrastructure : Cloud agnostique, multi-cloud Scaleway + AWS (~200K€/an total)",
            "Partenaires : INRIA (méthodologie), Gradient Health (données RWE), AMI Labs (évaluation indépendante)",
            "⚠ Clarification : « 100 milliards de relations » = encodeurs pré-entraînés externes (ChemBERTa, PubMedBERT, ESM-2), non entraînement from scratch",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q2: ETP par WP ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_title_text(slide, "Q2 — ETP prévus par WP et profils")
add_question_box(slide,
    "Quels sont les ETP prévus par WP et pour quels types de profils ?",
    "What are the planned FTEs per WP and for what types of profiles?")
add_wp_tag(slide, "Cross")
add_logo(slide)
add_footer(slide, page)
add_placeholder_box(slide, "RH / Direction", top=Inches(2.15))

# Empty FTE table as shapes
headers = ["WP", "IT (ETP)", "R&D (ETP)", "Médical (ETP)"]
wps = ["WP1", "WP2", "WP3", "WP4", "WP5", "WP6", "WP7"]
col_widths = [Inches(1.2), Inches(1.8), Inches(1.8), Inches(1.8)]
table_left = Inches(0.8)
table_top = Inches(2.85)
row_h = Inches(0.32)

# Header row
x = table_left
for i, h in enumerate(headers):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, col_widths[i], row_h)
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.color.rgb = WHITE
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    x += col_widths[i]

# Data rows (empty)
for r, wp in enumerate(wps):
    x = table_left
    for c in range(4):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, table_top + row_h * (r + 1),
            col_widths[c], row_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT_GREY
        box.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        box.line.width = Pt(0.5)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if c == 0:
            p.text = wp
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = DARK_NAVY
        else:
            p.text = "—"
            p.font.size = Pt(10)
            p.font.color.rgb = MID_GREY
        x += col_widths[c]

# Draft proposed FTEs
add_draft_box(slide,
    "Proposition ETP (basée sur les périmètres WP et plan de recrutement 2026) :\n"
    "• WP1 : IT 1.5 / R&D 2.0 / Méd 0.5  •  WP2 : IT 1.5 / R&D 2.0 / Méd 1.0\n"
    "• WP3 : IT 1.0 / R&D 2.0 / Méd 1.0  •  WP4 : IT 1.5 / R&D 1.5 / Méd 1.0\n"
    "• WP5 : IT 2.0 / R&D 1.5 / Méd 0.5  •  WP6 : IT 2.0 / R&D 3.0 / Méd 1.0\n"
    "• WP7 : IT 0.5 / R&D 0.5 / Méd 0.5  —  Total : ~27 ETP (aligné avec plan 47 postes dont 27 R&D)",
    "RH / Direction",
    top=Inches(5.2), height=Inches(1.6)
)
page += 1

# ─── Q3: Dissémination CIFRE ───
make_question_slide(
    q_num=3,
    short_title="Dissémination scientifique & CIFRE",
    q_fr="Dissémination scientifique : 4 doctorants CIFRE sont prévus (IA, pharmacologie et réglementaire). Pourront-ils publier ? Quelle articulation avec la stratégie PI (brevets et secret) ?",
    q_en="Will the 4 CIFRE doctoral students be able to publish? How does this align with the IP strategy (patents and trade secrets)?",
    wp_tag="Gouvernance",
    sections=[
        ("Expérience académique", [
            "Équipe dirigeante expérimentée : VS a dirigé 19 thèses dont 2 CIFRE (Sanofi, Servier)",
            "Pré-requis thèse CIFRE maîtrisés : équilibre découvertes scientifiques / applications industrielles",
        ]),
        ("Stratégie de publication", [
            "Les doctorants pourront publier — essentiel pour la validation de leur thèse",
            "Objectif : 5 publications annuelles pour s'établir comme acteur de référence international",
            "Positionnement sur la prédiction du bénéfice-risque de médicaments",
        ]),
        ("Protection IP", [
            "Brevets sur pipeline de gestion/indexation de données d'entraînement",
            "Secret industriel sur les architectures propriétaires et la Profiling Base",
            "Stratégie définie dès le sujet de thèse, réévaluée régulièrement",
            "Si obstacle à la publication anticipé (peu probable) → stratégie modifiée",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q4: LLM propriétaire ───
make_question_slide(
    q_num=4,
    short_title="LLM pour résumés et rapports",
    q_fr="Il est fait mention de LLM pour générer les résumés et rapports. S'agit-il d'un modèle propriétaire ? Quelle est la solution utilisée ?",
    q_en="LLMs are mentioned for generating summaries and reports. Is it a proprietary model? What solution is used?",
    wp_tag="Cross",
    sections=[
        ("Solution actuelle", [
            "Utilisation de Mistral, déployé sur serveurs GPU hébergés chez Scaleway",
        ]),
        ("Trois impératifs", [
            "Souveraineté des données : aucune donnée ne transite par des APIs tierces",
            "Souveraineté technologique : modèle open-source, pas de dépendance contractuelle",
            "Flexibilité : évaluation continue des nouveaux modèles (Llama, Gemma, Qwen)",
        ]),
        ("Architecture", [
            "Déploiement on-premise sur infrastructure Scaleway",
            "Capacité d'intégration rapide des modèles les plus performants sans contrainte",
        ]),
    ],
    speaker_notes="Note interne : Contrairement au dossier, nous utilisons actuellement des modèles OpenAI pour la summarization et les rapports. Ce point doit être clarifié avant la présentation.",
    page_num=page
)
page += 1

# ─── Q5: Transparence BDD & explicabilité ───
make_question_slide(
    q_num=5,
    short_title="Transparence des BDD et explicabilité",
    q_fr="Quelle transparence des BDD (publiques et privées) utilisées par la plateforme auprès des clients/utilisateurs finaux ? La traçabilité des analyses et des prédictions est un argument mis en avant, comment l'explicabilité se traduit-elle pour l'utilisateur ?",
    q_en="What transparency is provided regarding the databases used? How does explainability translate for the user?",
    wp_tag="Cross",
    sections=[
        ("Explicabilité des modèles (VS)", [
            "Essentielle dans la construction des modèles : vérification de cohérence",
            "Validation des relations entre mécanismes et outputs (efficacité / toxicité)",
            "Modèle explicable et validé = plus robuste qu'un modèle « black-box »",
        ]),
        ("Côté utilisateur final", [
            "Éléments mécanistiques explicables consultables à tout moment",
            "Essentiel = performance prédictive : VPP et VPN pour les applications d'intérêt",
        ]),
    ],
    draft_box=(
        "Bases de données utilisées :\n"
        "• Publiques : ChEMBL, PubMed, MEDLINE, ClinicalTrials.gov, FAERS, DisGeNET, PharmGKB, ClinVar, ToxCast, WITHDRAWN\n"
        "• Privées/partenaires : Sanofi (données internes), Cedars-Sinai (EDS), Mayo Clinic (EDS), ICM (cohortes)\n\n"
        "Explicabilité par couche :\n"
        "• SLM : provenance documentaire de chaque extraction\n"
        "• KG (WP5) : provenance au niveau de chaque arête, score de confiance, niveau de preuve\n"
        "• WP6 : SHAP values, attention weights, explications en langage naturel",
        "VS / Théo"
    ),
    page_num=page
)
page += 1

# ─── Q6: Trustworthy AI ───
make_question_slide(
    q_num=6,
    short_title="Problématiques de Trustworthy AI",
    q_fr="Comment sont adressées les problématiques de trustworthy AI ?",
    q_en="How are trustworthy AI issues addressed?",
    wp_tag="WP6",
    sections=[
        ("Explicabilité", [
            "SHAP, attention weights, explications en langage naturel (WP6)",
        ]),
        ("Traçabilité", [
            "MLflow pour versioning ; provenance au niveau des arêtes KG ; audit trail via API",
        ]),
        ("Quantification de l'incertitude", [
            "Variance inter-modèles (ensemble), OOD detection, intervalles de confiance explicites (T6.4)",
        ]),
        ("Conformité réglementaire", [
            "Architecture compatible ICH E2C(R2) et CIOMS XII ; alignement FDA, EMA, PMDA",
        ]),
        ("Positionnement éthique", [
            "Outil d'aide à la décision (usage pharma R&D), non dispositif médical",
        ]),
        ("Gouvernance des données", [
            "Accords institutionnels, dé-identification des données, approbations comités d'éthique (WP4)",
        ]),
    ],
    draft_box=(
        "Framework RAISE — mapping des piliers :\n"
        "• Accountable : MLflow audit trail, provenance KG, comité scientifique\n"
        "• Fair & Ethical : dé-identification, détection biais ethniques (WP3), biais de publication (WP2)\n"
        "• Robust & Safe : ensemble models, OOD detection, validation multi-niveaux, plan de mitigation par WP\n"
        "• Transparent & Explainable : SHAP, Mental Map (T6.5), explications en langage naturel\n"
        "• Eco-Responsible : optimisation architectures attention-free (T6.1), overhead <30%",
        "VS / Théo"
    ),
    page_num=page
)
page += 1

# ─── Q7: Biais World Model ───
make_question_slide(
    q_num=7,
    short_title="Détection et correction des biais",
    q_fr="Quelle est la stratégie pour détecter et corriger les biais dans le World Model ? Comment est-ce qu'elle s'intègre dans le planning du projet ?",
    q_en="What is the strategy for detecting and correcting biases in the World Model? How does it fit into the project timeline?",
    wp_tag="WP6",
    sections=[
        ("Biais identifiés", [
            "Biais de publication (WP2) — sous-représentation des résultats négatifs",
            "Biais de sélection d'efficacité (WP2)",
            "Biais de représentation ethnique dans les cohortes génomiques (WP3)",
            "Sous-déclaration dans les données RWE / FAERS (WP4)",
            "Biais des études publiées affectant la qualité des modèles (cross-WP)",
        ]),
        ("Stratégie de correction", [
            "Audit Knowledge Graph (WP5) : détection de contradictions entre relations",
            "Détection hors-distribution (OOD) : avertissements explicites en zone d'extrapolation",
            "Monitoring prospectif continu avec partenaires (T6.6)",
            "Détection de contradictions via KG comme contrôle qualité",
        ]),
    ],
    speaker_notes="Gap documenté : une stratégie formalisée et unifiée de correction des biais pour le World Model (WP6) n'est pas encore spécifiée dans le dossier — c'est un point ouvert identifié.",
    page_num=page
)
page += 1

# ─── Q8: Golden Standard ───
make_question_slide(
    q_num=8,
    short_title="Validation Golden Standard",
    q_fr="Concernant la validation des modèles avec un golden standard, de quel dataset s'agit-il ? (ex. p32)",
    q_en="Regarding model validation with a golden standard, which dataset is this?",
    wp_tag="Cross",
    sections=[
        ("Qualité Gold Standard", [
            "Données cliniques vérifiées par ARC qualifiés, critères diagnostiques selon référentiels récents",
            "Données manquantes <10%, stratégie claire d'imputation",
            "Contrôles qualité rigoureux sur données biologiques/imagerie",
        ]),
        ("NLP (WP2–WP4)", [
            "Datasets annotés manuellement avec cliniciens (F1 cible >0.85, complétude >0.70)",
            "Validation externe : Sanofi (96% couverture risque, 97% bénéfice), DisGeNET (92% top 30 biomarqueurs)",
        ]),
    ],
    two_col_sections=[
        ("Modèles prédictifs (WP1–WP4, WP6)", [
            "Panel de référence oncologie pulmonaire : succès (osimertinib, alectinib, lorlatinib) et échecs (buparlisib, idelalisib, selumetinib, vandetanib)",
            "Base WITHDRAWN : 578 médicaments retirés + molécules FDA boxed warning",
            "~4 500 molécules commercialisées (ChEMBL v35)",
        ]),
        ("World Model (WP6)", [
            "Validation rétrospective : ≥200 molécules profils B-R connus",
            "Validation prospective : monitoring chez partenaires (Sanofi, Cedars-Sinai, Mayo, ICM)",
            "Panel d'experts : ≥5 partenaires, cible ≥70% « cliniquement plausible »",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q9: Biais études rétrospectives ───
make_question_slide(
    q_num=9,
    short_title="Absence de biais — études rétrospectives",
    q_fr="Sur les études rétrospectives prévues, comment pensez-vous garantir l'absence de biais (ex. données relatives à la molécule commercialisée dans le(s) modèle(s) utilisé(s)) ?",
    q_en="For the planned retrospective studies, how do you intend to ensure the absence of bias?",
    wp_tag="Cross",
    sections=[
        ("Prévention fuite de données", [
            "Temporal split (hold-out 2022–2025), leave-drug-out cross-validation",
            "Nested cross-validation pour hyperparamètres ; audit conformité à chaque jalon",
            "WP1 : leave-one-scaffold-out pour généralisation chimique",
        ]),
        ("Biais de survie", [
            "Validation inclut explicitement WITHDRAWN (578 médicaments retirés) + FDA boxed warnings",
            "WP2 : inclut molécules ayant échoué en développement clinique",
        ]),
    ],
    two_col_sections=[
        ("Contrôle des facteurs confondants (WP4 RWE)", [
            "Analyse de sous-groupes : âge, comorbidités, ethnie déclarée",
            "Cross-validation stratifiée par source",
            "Ajustements statistiques pour biais de sélection par source RWE",
        ]),
        ("Barrière informationnelle", [
            "Panel de référence défini prospectivement, verrouillé avant développement WP1",
            "Jeu de validation WP6 (≥200 molécules) défini indépendamment",
            "Évaluation en aveugle pendant la phase de prédiction",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q10: Validation experts « cliniquement plausible » ───
make_question_slide(
    q_num=10,
    short_title="Méthodologie validation par experts",
    q_fr="Il est fait mention de la validation utilisateurs par des experts évaluateurs (de partenaires externes) jugeant les prédictions comme « cliniquement plausibles ». Pouvez-vous expliciter la méthodologie envisagée pour cette évaluation ?",
    q_en="Can you elaborate on the methodology for expert evaluator validation of predictions as 'clinically plausible'?",
    wp_tag="WP6",
    sections=[
        ("Approche knowledge-driven", [
            "Experts sélectionnés : profil physician-scientist, expertise physiopathologie / mécanismes / clinique",
            "Chaque prédiction = outcome (bénéfice/risque) + mécanisme(s) sous-jacent(s)",
        ]),
        ("Trois catégories de validation", [
            "1) Connexion mécanistique décrite dans la littérature → validation littérature",
            "2) Connexion non décrite mais compatible avec les connaissances existantes → nouveau mécanisme potentiel",
            "3) Connexion non décrite et incompatible → prédiction non validée",
        ]),
        ("Passage à l'échelle", [
            "Possibilité de compléter le travail expert par des LLMs pour valider un grand nombre d'hypothèses",
            "Objectif : mesures statistiquement significatives de puissance prédictive, sensibilité, spécificité",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q11: Monitoring ───
make_question_slide(
    q_num=11,
    short_title="Stratégie de monitoring",
    q_fr="Monitoring : quelle est la stratégie envisagée pour le World Model ? Qu'est-ce qui est mis en place actuellement sur la plateforme Trial Balancer ?",
    q_en="Monitoring: what is the strategy for the World Model? What is currently in place on Trial Balancer?",
    wp_tag="WP6",
    sections=[
        ("Trial Balancer — actuel", [
            "Réévaluation des modèles tous les 3–6 mois selon criticité",
            "Mise à jour continue des données : détection automatique de nouvelles publications",
            "SLM : F1 entre 80% et 95% ; modèle appraisal : précision 80–95%",
            "Monitoring infrastructure via Prometheus et Grafana",
        ]),
        ("World Model (BR-PREDICT) — planifié", [
            "Monitoring performance modèle : évaluation continue, retraining basé sur triggers (pas calendaire)",
            "Détection de drift : vérification automatique de distribution des features",
            "Maintenance KG (WP5) : mises à jour continues (nouveaux médicaments, cibles, EI, ontologies)",
            "Monitoring incertitude (T6.4) : signal intégré de dégradation",
            "MLflow : versionnage complet, rollback possible",
        ]),
    ],
    speaker_notes="Note interne — Certaines affirmations ne reflètent pas la réalité actuelle, notamment : existence des modèles WP6, les performances modèles décrites, et le monitoring infrastructure. Ces éléments sont des objectifs, pas l'état actuel. Le présentateur doit parler au futur et ne pas affirmer que ces capacités existent déjà.",
    page_num=page
)
page += 1

# ─── Q12: Déploiement on-premise ───
make_question_slide(
    q_num=12,
    short_title="Infrastructure déploiement on-premise",
    q_fr="Quelles sont les caractéristiques (stockage, ressources de calcul) des infrastructures nécessaires à un client pour le déploiement de la plateforme on-premise ?",
    q_en="What infrastructure characteristics are required by a client for on-premise deployment?",
    wp_tag="Cross",
    sections=[
        ("Architecture de déploiement", [
            "Docker (isolation) + Kubernetes (orchestration) + Ansible/Terraform (automatisation)",
            "Outils DevOps standards compatibles avec l'IT pharma existant",
            "Solution autonome post-installation ; accès ArcaScience compatible compliance client",
        ]),
        ("Calcul", [
            "Entraînement réalisé chez ArcaScience — client = inférence uniquement",
            "CPU : WP1, WP2 suffisent sur serveurs enterprise standards",
            "GPU : composants DL (WP3 GAT, WP4 TFT, WP6 ensemble) → un GPU enterprise (NVIDIA T4 ou A10)",
        ]),
    ],
    two_col_sections=[
        ("Stockage", [
            "Base structurée : capacités enterprise standards selon périmètre projet",
            "Stockage documentaire : Elasticsearch + vector embeddings pour recherche sémantique",
        ]),
        ("Sécurité & conformité", [
            "Chiffrement AES 256-bit, protocoles HIPAA-compatibles",
            "Données client indexées sur serveur du client uniquement",
            "Aucune donnée client ne transite par APIs externes ou services cloud",
        ]),
    ],
    speaker_notes="Note interne — Aucun de ces éléments ne reflète la réalité actuelle. Ce sont des capacités planifiées et des objectifs d'architecture. Le présentateur doit formuler ces points comme la cible de déploiement, pas comme l'état existant.",
    page_num=page
)
page += 1

# ─── Q13: Stratégie RH profils rares ───
make_question_slide(
    q_num=13,
    short_title="Stratégie RH — profils rares",
    q_fr="Le projet présente une stratégie de croissance RH ambitieuse. Quelle est la stratégie pour s'assurer de trouver les profils adéquats, potentiellement rares sur le marché ?",
    q_en="What is the strategy for ensuring that adequate profiles, potentially rare on the market, can be found?",
    wp_tag="Gouvernance",
    sections=[],
    placeholder="RH / Direction",
    draft_box=(
        "Proposition de réponse :\n"
        "• Pipeline CIFRE : 4 doctorants en IA, pharmacologie et réglementaire — vivier académique directement intégré au projet\n"
        "• Co-supervision INRIA : accès au réseau de chercheurs seniors en ML/IA pour les profils R&D les plus pointus\n"
        "• Sourcing international : partenariats Sanofi, Cedars-Sinai, Mayo Clinic comme canaux de recrutement (physician-scientists, data scientists cliniques)\n"
        "• Plan de recrutement 2026 : 47 postes dont 27 R&D — recrutement échelonné aligné sur les jalons WP\n"
        "• Packages compétitifs : equity + participation au projet I-demo comme levier d'attractivité\n"
        "• Équipe fondatrice expérimentée : VS (19 thèses dirigées), expertise en encadrement scientifique",
        "RH / Direction"
    ),
    page_num=page
)
page += 1

# ─── Q14: Scénarios alternatifs ───
make_question_slide(
    q_num=14,
    short_title="Scénarios alternatifs — défaillance WP",
    q_fr="Y a-t-il des scénarios alternatifs prévus en cas de défaillance d'un WP (difficulté d'accès à une BDD, absence de validation d'un modèle…) ?",
    q_en="Are there alternative scenarios in case of WP failure (database access issues, model validation failure)?",
    wp_tag="Cross",
    sections=[
        ("Accès aux données", [
            "Anticipation : jusqu'à 50% des BDD non exploitables (qualité, contrats, délais)",
            "Stratégie : multiplication des sources et partenaires — aucun n'est indispensable",
            "WP4 (RWE) : risque plus élevé mais non bloquant — FAERS public + WP1-3 suffisants pour le modèle",
            "WP1-3 : sources principalement publiques (ChEMBL, ToxCast, PharmGKB, ClinVar…) — risque faible",
        ]),
        ("Validation insuffisante d'un modèle", [
            "Validation = notion multi-niveaux : 1) in silico, 2) basée sur les connaissances, 3) expérimentale",
            "Validation quantitative et par tâche — un modèle peut être validé partiellement",
            "Standards de validation multi-points, qualitative et quantitative appliqués systématiquement",
            "Si validation insuffisante → analyse des causes → 1) boucle rétroactive, 2) ajout de données ciblées",
        ]),
        ("Architecture résiliente", [
            "WP1-4 indépendants : chacun délivre une valeur standalone",
            "WP6 World Model : gère les modalités manquantes via learned mask tokens (mode dégradé robuste)",
            "WP6 = axe exploratoire — recherche de frontière, validée par expériences focalisées avant engagement complet",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q15: INRIA sous-traitance ───
make_question_slide(
    q_num=15,
    short_title="Sous-traitance INRIA et PI",
    q_fr="Quelle est la nature de la sous-traitance de l'INRIA et quels sont les profils impliqués du côté de ce prestataire ? Y a-t-il un partage de la PI ?",
    q_en="What is the nature of INRIA's subcontracting? What profiles are involved? Is there IP sharing?",
    wp_tag="WP5",
    sections=[],
    placeholder="Direction / Cabinet Carrel",
    draft_box=(
        "Proposition de réponse :\n"
        "• Nature : support méthodologique sur WP5 (Knowledge Graph) et WP6 (World Model) — expertise en graphes de connaissances, raisonnement causal, modélisation ML\n"
        "• Profils INRIA : chercheurs seniors en ML/IA, spécialistes graphes et ontologies\n"
        "• Publications conjointes : prévues dans le cadre des thèses CIFRE co-supervisées\n"
        "• PI — Contrat standard INRIA :\n"
        "  - Background IP : retenue par chaque partie\n"
        "  - Foreground IP : co-propriété avec droits d'exploitation exclusifs pour ArcaScience dans le domaine pharma B-R\n"
        "  - Publications : soumises à un délai de confidentialité (typiquement 6 mois) pour permettre le dépôt de brevets",
        "Direction / Cabinet Carrel"
    ),
    page_num=page
)
page += 1

# ─── Q16: Traçabilité et qualité ───
make_question_slide(
    q_num=16,
    short_title="Management traçabilité et qualité",
    q_fr="Quelle est la stratégie de management de la traçabilité et de la qualité des développements prévus dans les WP ? Quelles méthodes et outils seront mis en place ?",
    q_en="What is the management strategy for traceability and quality of the developments planned in the WPs?",
    wp_tag="Cross",
    sections=[
        ("Cadre de contrôle qualité", [
            "Contrôle qualité à toutes les étapes : cohortes patients, données, traitement, intégration, modèles, performances",
            "SOPs préparées à l'avance et suivies rigoureusement",
            "SOPs transmises à tous les partenaires pour garantir homogénéité des procédures",
        ]),
        ("Outils", [
            "MLflow : versioning modèles, tracking expériences, stockage artefacts",
            "SonarQube : quality gate sur chaque mise en production",
            "Prometheus / Grafana : monitoring infrastructure et disponibilité",
            "Audit trail via API sécurisée",
        ]),
        ("Attention particulière aux biais", [
            "Détection et mitigation systématiques à chaque étape",
            "Objectif : modèle basé sur éléments cliniques et biologiques robustes",
            "Diminution du bruit et des biais inhérents aux données biomédicales",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q17: T7.3 Discussions pharma s'arrêtent à M24 ───
make_question_slide(
    q_num=17,
    short_title="Discussions partenariales pharma (T7.3)",
    q_fr="Pourquoi les discussions partenariales avec l'industrie pharmaceutique (T7.3) s'arrêtent à M24 sur le diagramme de Gantt ? Quelle est la portée de ces discussions ?",
    q_en="Why do partnership discussions with pharma (T7.3) stop at M24? What is their scope?",
    wp_tag="WP7",
    sections=[],
    placeholder="Direction / Stratégie commerciale",
    draft_box=(
        "Proposition de réponse :\n"
        "• Portée des discussions T7.3 : accès aux BDD privées, accès aux datasets de validation, accords de co-développement, contrats early adopter\n"
        "• Le Gantt ne montre que la tâche formelle « négociation partenariale » — les discussions continuent au-delà de M24 sous le volet commercial de WP7\n"
        "• Après M24, les partenariats signés entrent en phase d'exécution (accès données, validation, co-développement)\n"
        "• Pipeline commercial actif : 5-7 nouveaux logos ciblés par trimestre, expansion des comptes existants\n"
        "• Modèle de pricing disease-specific : BRA Essentials (75-100K€/an), Professional (125-175K€/an), Enterprise (200-300K€/an)",
        "Direction / Stratégie commerciale"
    ),
    page_num=page
)
page += 1

# ─── Q18: KG publics — concurrence et interopérabilité ───
make_question_slide(
    q_num=18,
    short_title="KG publics et interopérabilité",
    q_fr="WP5 : les graphes publics de référence sont-ils utilisés par des concurrents ? Pourquoi une interopérabilité avec ces KG n'est-elle pas envisagée ?",
    q_en="Are public reference graphs used by competitors? Why is interoperability with these KGs not considered?",
    wp_tag="WP5",
    sections=[
        ("Positionnement vs KG publics", [
            "DRKG (5.9M arêtes non pondérées), PrimeKG (4M relations sans statut causal) cités comme baseline",
            "BR-PREDICT KG : >95% arêtes annotées avec scores de confiance, provenance, niveau de preuve",
            "Chaque arête typée : causal-validé, directionnel-inféré, ou associatif — seuls les deux premiers alimentent le SCM (WP6)",
        ]),
        ("Choix délibéré de non-import", [
            "KG publics sans statut causal, scores de confiance, ni provenance → incompatibilité qualité",
            "Import contaminerait le KG avec relations non qualifiées → dégradation auditabilité réglementaire",
            "Exigences ICH E2C(R2), CIOMS XII imposent traçabilité complète",
        ]),
        ("Interopérabilité effective", [
            "Stockage RDF parallèle pour interopérabilité SPARQL",
            "Alignement sur ontologies standards : MedDRA, ChEBI, Gene Ontology, Reactome, Disease Ontology",
            "Architecture interopérable au niveau ontologique — ne consomme pas les KG publics comme input",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q19: AMI Labs en fin de projet ───
make_question_slide(
    q_num=19,
    short_title="AMI Labs — évaluation précoce",
    q_fr="AMI Labs intervient en tant que prestataire en fin de projet pour définir la méthodologie d'évaluation. Quelles sont les motivations de ne pas intégrer l'évaluation dès les premières phases de développement de la solution ?",
    q_en="Why is AMI Labs' evaluation methodology not integrated from the earliest development phases?",
    wp_tag="WP6",
    sections=[
        ("Raison pragmatique", [
            "Le World Model n'existe pas avant M18–M30 → rien à évaluer de façon intégrée avant",
            "AMI Labs intervient quand le système est suffisamment mature pour une évaluation méthodologique indépendante",
        ]),
        ("Évaluation per-WP existante dès le début", [
            "Chaque WP (1–4) dispose de son propre protocole de validation avec KPIs définis",
            "Validation continue tout au long du développement (F1, AUC, calibration, etc.)",
        ]),
        ("Amélioration identifiée", [
            "Une consultation méthodologique précoce avec AMI Labs sur les protocoles per-WP aurait été justifiable",
            "Ce point sera intégré plus tôt dans la réalisation du projet",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q20: Structure capitalistique ───
make_question_slide(
    q_num=20,
    short_title="Structure capitalistique et juridique",
    q_fr="Une opération affectant la structure juridique et/ou capitalistique de votre structure est-elle envisagée dans les deux ans à venir ?",
    q_en="Is any operation affecting the legal and/or capital structure planned within the next two years?",
    wp_tag="Gouvernance",
    sections=[],
    placeholder="Direction / Cabinet Carrel",
    draft_box=(
        "Proposition de réponse :\n"
        "• Levée de fonds prévue dans les 12-18 prochains mois pour soutenir l'exécution du plan I-demo\n"
        "• Structure juridique stable — SAS de droit français\n"
        "• Pas de fusion, scission ou changement de forme juridique envisagé\n"
        "• Accompagnement juridique : Cabinet Carrel (fondé par ex-directrice Inserm Transfert)",
        "Direction / CEO"
    ),
    page_num=page
)
page += 1

# ─── Q21: Levée de fonds ───
make_question_slide(
    q_num=21,
    short_title="Levée de fonds",
    q_fr="Votre structure envisage-t-elle une levée de fonds dans les deux ans à venir ? Si oui, pouvez-vous nous indiquer l'identité des investisseurs/fonds d'investissement, les montants envisagés et le calendrier ?",
    q_en="Is a fundraising round planned? If so, investor identity, amounts, and timeline?",
    wp_tag="Gouvernance",
    sections=[],
    placeholder="Direction / CEO",
    draft_box=(
        "Proposition de réponse :\n"
        "• Levée de fonds en cours de structuration pour H2 2026 – H1 2027\n"
        "• Montant, identité des investisseurs et calendrier : à compléter par la direction\n"
        "• Objectif : financer l'accélération commerciale (O3 — $500K ARR run-rate) et l'exécution des WP6-7",
        "CEO"
    ),
    page_num=page
)
page += 1

# ─── Q22: Financements publics ───
make_question_slide(
    q_num=22,
    short_title="Sources publiques de financement",
    q_fr="Êtes-vous déjà bénéficiaire, ou en cours de demande, directement ou indirectement au titre de ce projet d'autres financements (financements nationaux, régionaux, aides européennes, etc.) ?",
    q_en="Are you already receiving or applying for other public funding for this project?",
    wp_tag="Gouvernance",
    sections=[],
    placeholder="Direction / DAF",
    draft_box=(
        "Proposition de réponse :\n"
        "• Détailler les financements existants (CIR, JEI, subventions régionales, etc.)\n"
        "• Lister les demandes en cours\n"
        "• Confirmer la conformité avec les règles de cumul d'aides Bpifrance",
        "Direction / DAF"
    ),
    page_num=page
)
page += 1

# ─── Q23: Sous-traitance contrats ───
make_question_slide(
    q_num=23,
    short_title="Accords de sous-traitance",
    q_fr="Des sous-traitances dans le cadre du projet sont envisagées, pouvez-vous nous fournir les accords/modèles d'accord encadrant ces relations de sous-traitance ? Si un sous-traitant est défaillant, avez-vous identifié des alternatives ?",
    q_en="Can you provide subcontracting agreements? Have you identified alternatives if a subcontractor fails?",
    wp_tag="Gouvernance",
    sections=[
        ("Stratégie de mitigation", [
            "Gradient Health : écosystème large de partenaires données — multiplication des sources pour compenser les défections",
            "Cedars-Sinai / Mayo Clinic : fournisseurs de données et sites de validation",
            "ICM : cohortes génomiques et cliniques",
            "Aucun partenaire unique n'est indispensable — couverture redondante par design",
        ]),
        ("Accords contractuels", [
            "Modèles d'accord à fournir au comité (SOPs et templates contractuels existants)",
            "Accompagnement juridique : Cabinet Carrel (fondé par ex-directrice Inserm Transfert)",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q24: Veille réglementaire ───
make_question_slide(
    q_num=24,
    short_title="Veille réglementaire",
    q_fr="Comment est faite la veille réglementaire au sein de votre structure ?",
    q_en="How is regulatory monitoring conducted within your organization?",
    wp_tag="Gouvernance",
    sections=[
        ("Organisation interne", [
            "1 Q&A GxP (Charbel / Clarisse) — conformité pharmaceutique et réglementaire",
            "1 Q&A IT (Théo) — conformité technique et cybersécurité",
            "1 Comité scientifique — veille scientifique et méthodologique",
        ]),
        ("Périmètre de veille", [
            "Réglementation pharmaceutique : ICH, EMA, FDA, PMDA",
            "Réglementation IA : EU AI Act, guidelines nationales",
            "Protection des données : RGPD, HIPAA",
            "Dispositifs médicaux : MDR/IVDR (veille — non applicable dans le périmètre actuel)",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q25: Verrous juridiques ───
make_question_slide(
    q_num=25,
    short_title="Verrous juridiques",
    q_fr="Comment et par qui êtes-vous accompagnés sur le plan juridique et réglementaire ? Est-ce que vous avez relevé des verrous juridiques pour mener à bien ce projet ?",
    q_en="Who provides legal and regulatory support? Have you identified any legal barriers?",
    wp_tag="Gouvernance",
    sections=[
        ("Accompagnement juridique", [
            "Cabinet Carrel — fondé par l'ex-directrice d'Inserm Transfert",
            "Expertise en propriété intellectuelle, contrats de recherche collaborative, droit de la santé",
        ]),
        ("Verrous juridiques", [
            "Aucun verrou juridique identifié à ce stade",
            "Les verrous ont été levés et intégrés au développement d'ArcaScience au fil des années",
            "Veille continue pour anticiper les évolutions réglementaires (EU AI Act, EHDS, etc.)",
        ]),
    ],
    page_num=page
)
page += 1

# ─── Q26: Retour clients SaaS ───
make_question_slide(
    q_num=26,
    short_title="Retour clients mode SaaS",
    q_fr="Pourriez-vous nous donner le retour de vos premiers clients sur le mode SaaS ?",
    q_en="Could you give us feedback from your first SaaS clients?",
    wp_tag="WP7",
    sections=[],
    placeholder="Direction commerciale / CEO",
    draft_box=(
        "Proposition de réponse :\n"
        "• Plateforme BRA en phase de déploiement auprès des premiers utilisateurs\n"
        "• Objectif Q en cours : 10-50 utilisateurs de 1-5 entreprises (KR 1.1)\n"
        "• Pricing SaaS structuré : BRA Essentials (75-100K€/an), Professional (125-175K€/an), Enterprise (200-300K€/an)\n"
        "• Retours qualitatifs à collecter et présenter (feedback pipeline, NPS, cas d'usage validés)",
        "CEO / Direction commerciale"
    ),
    page_num=page
)
page += 1

# ─── Q27: Contrat ICON ───
make_question_slide(
    q_num=27,
    short_title="Contrat ICON",
    q_fr="Où en est le contrat avec ICON ?",
    q_en="What is the status of the contract with ICON?",
    wp_tag="WP7",
    sections=[],
    placeholder="CEO / Direction commerciale",
    draft_box=(
        "Proposition de réponse :\n"
        "• Statut actuel du contrat ICON à détailler\n"
        "• Périmètre de la collaboration, montant, jalons\n"
        "• Lien avec la stratégie go-to-market et le pipeline commercial",
        "CEO"
    ),
    page_num=page
)
page += 1

# ─── Q28: Contrat Vidal / Synapse ───
make_question_slide(
    q_num=28,
    short_title="Contrat Vidal et module Synapse",
    q_fr="Pourriez-vous donner plus d'informations sur votre contrat avec Vidal et comment ce nouveau module peut ou pourrait s'articuler avec les modules d'analyse de risques d'interactions médicamenteuses comme celui de Synapse ?",
    q_en="More details on the Vidal contract and articulation with drug interaction risk analysis modules like Synapse?",
    wp_tag="WP7",
    sections=[
        ("Positionnement différenciant", [
            "Technologie Synapse/Vidal : compréhension des leaflets et du travail du G-TIAM",
            "ArcaScience : en amont — identification des cas d'interactions médicamenteuses issus de la vie réelle et de la littérature",
            "Complémentarité : ArcaScience enrichit les signaux en amont, Vidal les diffuse en aval",
        ]),
    ],
    placeholder="CEO / Direction commerciale",
    draft_box=(
        "Proposition de réponse :\n"
        "• Détailler le périmètre du contrat Vidal (nature, montant, durée)\n"
        "• Expliquer le module spécifique développé ou en cours de développement\n"
        "• Articuler avec la stratégie d'interactions médicamenteuses de BR-PREDICT (T6.3 SCM)",
        "CEO / Direction commerciale"
    ),
    page_num=page
)
page += 1

# ─── Q29: Organisation sales ───
make_question_slide(
    q_num=29,
    short_title="Identification et contact clients",
    q_fr="Comment êtes-vous structurés pour l'identification et la prise de contact auprès des clients potentiels ?",
    q_en="How are you structured for identifying and contacting potential clients?",
    wp_tag="WP7",
    sections=[],
    placeholder="CEO / Direction commerciale",
    draft_box=(
        "Proposition de réponse :\n"
        "• Organisation sales structurée avec pipeline discipliné\n"
        "• Target : 10 outbounds + 5 warm intros par quinzaine (Weeks 1-2 du quarter)\n"
        "• Pipeline coverage target : 3x du target revenue ($1.5M)\n"
        "• Motion disease-specific-first : les versions thérapeutiques compressent les cycles de vente de 40-60%\n"
        "• Win rate cible : 25-35% sur opportunités qualifiées\n"
        "• Cycle de vente moyen : <60 jours\n"
        "• Profils cibles : Top-20 pharma, mid-market biotech, CROs",
        "CEO"
    ),
    page_num=page
)
page += 1

# ─── Q30: Positionnement NSCLC ───
make_question_slide(
    q_num=30,
    short_title="Choix du NSCLC et indications futures",
    q_fr="Pourquoi avoir choisi le NSCLC comme première indication ? Quelles seraient les indications suivantes ? Comment allez-vous sourcer les données pour les 10% spécifiques par indication ?",
    q_en="Why NSCLC first? What are the next indications? How will you source the 10% indication-specific data?",
    wp_tag="Stratégie",
    sections=[
        ("Choix NSCLC", [
            "Stratégie orientée cancer : non-small cell lung cancer comme socle",
            "Pertinence scientifique et richesse de données maximales",
            "Création d'une base de données critique transférable",
        ]),
        ("Indications suivantes", [
            "Extension vers l'immuno-inflammation en fonction de la disponibilité des données",
            "Transfer learning entre aires thérapeutiques",
            "Similarité des pathologies facilite le transfert de modèles",
            "Architecture 90/10 : ~10% de calibration spécifique par nouvelle indication",
        ]),
        ("Sourcing des 10% spécifiques", [
            "1) Moissonnage : données open structurées + non-structurées (PubMed, ChEMBL, ClinicalTrials.gov, FAERS)",
            "2) Partenariats privés : Gradient Health, Cedars-Sinai, Mayo Clinic",
            "3) Partenariats publics : ICM, Institut Bergonié",
        ]),
    ],
    page_num=page
)
page += 1

# ═══════════════════════════════════════════════════════════════
# CLOSING SLIDE: Prochaines étapes
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_logo(slide)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Prochaines étapes"
run.font.size = Pt(32)
run.font.bold = True
run.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

bullets_data = [
    "1. Finaliser les réponses DRAFT flaguées — chaque owner valide son périmètre",
    "2. Valider les aspects juridiques et PI avec Cabinet Carrel + revue RH",
    "3. Renvoyer la version consolidée au comité d'évaluation Bpifrance",
]
txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(9.3), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, b in enumerate(bullets_data):
    p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
    run = p.add_run()
    run.text = b
    run.font.size = Pt(18)
    run.font.color.rgb = CYAN
    p.space_before = Pt(16)

# Accent line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(4.5), Inches(2.6),
    Inches(4.3), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = AMBER
line.line.fill.background()

add_footer(slide, page)
page += 1

# ═══════════════════════════════════════════════════════════════
# APPENDIX: DRAFT / À compléter items
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_title_text(slide, "Annexe — Items DRAFT & À compléter")
add_logo(slide)
add_footer(slide, page)

draft_items = [
    ("Q2", "ETP par WP — table à compléter", "RH / Direction"),
    ("Q5", "Liste BDD publiques/privées + explicabilité par couche", "VS / Théo"),
    ("Q6", "Framework RAISE Trustworthy AI", "VS / Théo"),
    ("Q13", "Stratégie RH profils rares", "RH / Direction"),
    ("Q15", "Sous-traitance INRIA — nature, profils, PI", "Direction / Cabinet Carrel"),
    ("Q17", "T7.3 discussions pharma après M24", "Direction / Stratégie"),
    ("Q20", "Structure capitalistique", "Direction / CEO"),
    ("Q21", "Levée de fonds — montant, investisseurs, calendrier", "CEO"),
    ("Q22", "Sources publiques de financement", "Direction / DAF"),
    ("Q26", "Retour clients SaaS", "CEO / Direction commerciale"),
    ("Q27", "Contrat ICON", "CEO"),
    ("Q28", "Contrat Vidal / Synapse", "CEO / Direction commerciale"),
    ("Q29", "Organisation sales", "CEO"),
]

# Table
col_w = [Inches(1.0), Inches(6.5), Inches(3.5)]
y = Inches(1.1)
# Header
x = Inches(1.0)
for i, h in enumerate(["Question", "Item à compléter", "Owner"]):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w[i], Inches(0.35))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.color.rgb = WHITE
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    x += col_w[i]

y += Inches(0.35)
for r, (q, item, owner) in enumerate(draft_items):
    x = Inches(1.0)
    vals = [q, item, owner]
    for c in range(3):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w[c], Inches(0.35))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT_GREY
        box.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        box.line.width = Pt(0.5)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = vals[c]
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_GREY if c > 0 else ORANGE
        p.font.bold = (c == 0)
        x += col_w[c]
    y += Inches(0.35)

# ─── Save ───
output_path = os.path.join(os.path.dirname(__file__), "ArcaScience_I-demo_QA_Committee_Deck_v1.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
